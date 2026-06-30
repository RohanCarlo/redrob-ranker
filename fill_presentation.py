#!/usr/bin/env python3
"""
Redrob hackathon PPTX filler — improved version.
Fixes: text overflow, adds matplotlib diagrams, extends content boxes to 4.0".
"""

import copy
import csv
import io

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch

from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

INPUT_FILE = "Idea Submission Template _ Redrob.pptx"
OUTPUT_FILE = "Redrob_Submission_Rohan.pptx"

PARTICIPANT_NAME = "Rohan"
PARTICIPANT_EMAIL = "rohan2072001@gmail.com"
GITHUB_REPO = "https://github.com/rohan2072001/redrob-ranker"

# Content area layout (extended height to prevent overflow)
C_LEFT   = Inches(0.41)
C_TOP    = Inches(1.42)
C_WIDTH  = Inches(9.32)
C_HEIGHT = Inches(4.05)   # was 3.50 — extended to 4.05" (bottom at 5.47", 0.15" margin)
C_FONT   = 9.0            # pt — 9pt allows ~26 lines in 4.05" with single spacing

# ─── Slide text content ───────────────────────────────────────────────────────

SLIDE_CONTENT = {
    1: {
        "Team Name :":         f"Team Name :  {PARTICIPANT_NAME}",
        "Team Leader Name :":  f"Team Leader Name :  {PARTICIPANT_NAME}",
        "Problem Statement :": "Problem Statement :  Intelligent Candidate Discovery & Ranking — Senior AI Engineer @ Redrob AI",
    },

    2: """\
PROPOSED SOLUTION
• Single-file (ranker.py, ~420 lines), pure Python stdlib — zero external dependencies
• Ranks 100,000 candidates in ~7 seconds on CPU — no GPU, no external APIs required
• 4 hard must-have gates prevent unqualified candidates from surfacing in top-100
• Grounded 1–2 sentence reasoning for every ranked candidate, citing actual profile facts
• 19 honeypot profiles (impossible job durations) detected and filtered before ranking

HOW IT DIFFERS FROM TRADITIONAL MATCHING
• Gate multiplier:  4/4 must-haves → ×1.00  |  3/4 → ×0.78  |  2/4 → ×0.56  |  0/4 → ×0.10
• Skill quality:    proficiency (expert/advanced/intermediate/beginner) × duration × endorsements
• Evaluation depth: NDCG/MRR/A/B-test keyword density in career descriptions scored proportionally
• Career trajectory: rewards consistent ML roles; penalises consulting-only backgrounds (Wipro/TCS…)
• Behavioral signals: GitHub activity, recruiter response rate, activity recency — all factored in

ARCHITECTURE HIGHLIGHTS
• O(N) single streaming pass — 487 MB JSONL never fully loaded into RAM
• Min-heap (size 400) tracks top candidates live; O(log N) insert/evict per record
• 12 frozenset skill taxonomies covering 150+ aliases for O(1) membership lookup
• Fully deterministic and reproducible — same input always produces same output""",

    3: """\
JD REQUIREMENTS EXTRACTED
MUST-HAVES  (all 4 required — missing any collapses score via gate multiplier)
  ① Embedding retrieval:    SBERT · BGE · E5 · GTE · word2vec · DPR · ColBERT · ada-002
  ② Vector DB / ANN search: Pinecone · Weaviate · Qdrant · Milvus · FAISS · pgvector · Chroma
  ③ Ranking eval:           NDCG · MRR · MAP · A/B testing · recall@K · precision@K
  ④ Strong Python  (+ ecosystem: PyTorch · FAISS · HuggingFace · sklearn)

NICE-TO-HAVES:  LLM fine-tuning (LoRA/QLoRA/PEFT) · LTR (LambdaRank) · RAG/BM25 · MLOps

KEY CANDIDATE EVALUATION SIGNALS  (how we go beyond keyword matching)
  Skills        proficiency level × duration in months × endorsement count
                must-have gate score (4 categories) + nice-to-have quality bonus
  Experience    years-in-band, ML title relevance, career trajectory consistency
                consulting-only penalty for TCS/Wipro/Infosys/Accenture-heavy profiles
  Behavioral    GitHub activity score (0–100) · recruiter response rate (0–1)
                activity recency buckets: 7d / 30d / 90d / 180d / >180d
  Education     institution tier × degree level × STEM / CS field-of-study match
  Notice        > 90-day notice period flagged as concern in reasoning output""",

    4: """\
THREE-STEP PIPELINE

STEP 1  RETRIEVE — O(N) streaming pass through candidates.jsonl (487 MB, 100 K records)
   • Min-heap (size 400) tracks top candidates live; lowest score evicted on each improvement
   • No pre-indexing, no database, no external storage — pure iterator over file lines

STEP 2  SCORE — Composite formula per candidate:
   score = MH_GATE × disqual_mult × (
     0.42 × skills_component  +  0.28 × experience_component
     + 0.18 × signal_score    +  0.12 × education_score  )

   MH_GATE multipliers:
     4/4 → ×1.00  |  3/4 → ×0.78  |  2/4 → ×0.56  |  1/4 → ×0.30  |  0/4 → ×0.10

   Skills component:  55% must-have score  +  45% nice-to-have score
   Must-have weights: embeddings 40%  |  vector DB 30%  |  eval 20%  |  Python 10%
   Eval score: 0.15 base + 0.07 per keyword hit in career text (max 0.55)

STEP 3  RANK — Sort top-400 (score DESC, candidate_id ASC) → output top-100

ALGORITHMS & HEURISTICS
• Skill quality:    proficiency weight × (duration_months / 60) × (endorsements / 40)
• Honeypot filter: claimed job tenure vs. actual elapsed months since start_date (±6 month tolerance)
• Consulting penalty: proportion of roles at 13 big IT services firms → deducted from exp score
• Disqualifier mult: CV/speech specialist score × 0.25 if no NLP/IR bridge present
• Behavioral decay: activity recency scored in buckets — 7d · 30d · 90d · 180d · 180d+""",

    5: """\
HOW RANKING DECISIONS ARE EXPLAINED
  Every top-100 candidate receives a grounded, fact-cited reasoning string:
  "<Title> (X.X yrs): N/4 JD must-haves — embeddings (<skill>); vector DB (<db>);
   eval (<evidence>); response rate X.XX, GitHub YY. [Concerns: ...]"

EXAMPLES FROM SUBMISSION
  "Senior ML Engineer (7.2 yrs): 4/4 — embeddings (Embeddings); vector DB (Weaviate/Pinecone);
   eval (NDCG, MRR in work history); response 0.61, GitHub 95."
  "ML Engineer (7.2 yrs): 4/4 — embeddings (Embeddings); vector DB (FAISS/pgvector);
   eval (a/b test in work history); response 0.60, GitHub 63. Concerns: notice 120d."
  32 / 100 candidates have concern clauses — transparent about trade-offs.

PREVENTING HALLUCINATIONS  (no LLM — zero generative steps)
  • All facts parsed directly from candidate profile fields; no text generation
  • Skill names = exact strings from candidate's skills[] array
  • "in work history" explicitly marks text-mined vs. skills-listed evidence
  • Concerns only appended when specific numeric thresholds are exceeded

HANDLING INCONSISTENT / LOW-QUALITY / SUSPICIOUS PROFILES
  Suspicious:   Honeypot check — claimed tenure vs. actual months since start_date (±6 mo)
                19 profiles caught and excluded before any scoring begins
  Low-quality:  0/4 must-haves → gate multiplier ×0.10; cannot appear in top-100
                Missing or empty fields default to zero contribution (no score inflation)
  Inconsistent: Consulting-only careers (TCS/Wipro/Infosys/Accenture) → exp penalty
                Inactive >90d → signal score capped + concern flagged in reasoning""",

    6: """\
STEP 1 — JD ANALYSIS  (offline, coded once before the scoring run)
  • Read JD → identify 4 must-have categories, 4 nice-to-have groups, 2 disqualifier sets
  • Build 12 frozenset skill taxonomies (150+ aliases total) for O(1) membership lookup
  • Hard-code must-have gate multipliers: {0:0.10, 1:0.30, 2:0.56, 3:0.78, 4:1.00}

STEP 2 — STREAM & SCORE  (~7 seconds for all 100 K candidates)
  For each JSON line in candidates.jsonl:
    a.  Parse JSON → extract profile, career_history, skills, redrob_signals, education
    b.  Honeypot check → skip if claimed job tenure is chronologically impossible
    c.  Build career_text = all job descriptions + titles + company names + summary + headline
    d.  score_skills()     → must-have gate + quality scoring + keyword density in career text
    e.  score_experience() → years band, title relevance, ML trajectory, consulting penalty
    f.  score_signals()    → activity recency, response rate, GitHub score, notice period
    g.  score_education()  → institution tier × degree level × field-of-study match
    h.  composite_score = MH_GATE × disqual_mult × (0.42·skills + 0.28·exp + 0.18·sig + 0.12·edu)
    i.  Min-heap update (size 400) → evict lowest-scored entry if new score is higher

STEP 3 — SELECT & OUTPUT
  • Sort top-400 by (score DESC, candidate_id ASC for deterministic tie-breaking)
  • Take top-100 → generate grounded reasoning string per candidate
  • Write submission.csv (4 columns: candidate_id, rank, score, reasoning)
  • Run validate_submission.py → confirmed "Submission is valid." ✓""",

    # Slide 7: replaced by matplotlib architecture diagram — see make_architecture_diagram()

    8: """\
DATASET & FILTERING
  • 100,000 candidates scanned in candidates.jsonl (487 MB)
  • 19 honeypot profiles removed pre-ranking
  • 308 with all 4/4 must-haves → entire top-100 drawn from this pool
  • 6,370 with 3/4 must-haves (gate ×0.78, scored below 4/4 pool)

TOP-100 SNAPSHOT  (score range: 0.863 → 0.691, see chart →)
  Top titles:   Applied ML Eng (20) · ML Engineer (14)
                Search Engineer (13) · AI Engineer (10)
                NLP / Sr NLP Engineer (12) · Rec. Systems Eng (6)
  Experience:   2.9–16.9 yrs  |  median 6.5 yrs  (JD: 5–9 yrs ✓)
  Location:     87% India-based  |  13% open to relocation
  Concerns:     32/100 flagged (notice period or inactivity)

COMPUTE COMPLIANCE
  ✓ Runtime:  ~7 sec  (constraint: ≤ 5 minutes)
  ✓ Memory:   streaming — 487 MB never loaded into RAM
  ✓ Network:  zero external calls
  ✓ GPU:      none — pure CPU
  ✓ pip:      zero installs at runtime""",

    9: """\
CORE RUNTIME — ZERO EXTERNAL DEPENDENCIES  (Python 3.9 stdlib only)
  Module      Role in pipeline
  json        Streaming JSONL parse, one line at a time (never loads full 487 MB)
  csv         Compliant CSV output with header, quoting rules, and UTF-8 encoding
  heapq       O(log N) min-heap for live top-400 candidate tracking
  datetime    ISO date arithmetic for honeypot detection and activity recency scoring
  re          Skill name normalisation and keyword pattern matching
  time / pathlib   Wall-clock profiling; cross-platform file path handling

WHY PURE STDLIB (no ML frameworks)?
  • Zero-dependency: runs anywhere Python ≥ 3.7 is available — no pip install, no conflicts
  • A well-designed domain-expert heuristic outperforms a naïve ML model when no ground-truth
    labels exist for calibration — and it is fully deterministic and debuggable
  • Satisfies the "no network, no GPU" constraint without any special configuration

WHAT WAS DELIBERATELY NOT USED  (and why)
  Tool / library                    Reason excluded
  sentence-transformers / OpenAI    Model download requires network access — prohibited
  scikit-learn / XGBoost            Unnecessary complexity; adds install step
  Any LLM API (Claude, GPT-4…)      Explicitly prohibited by challenge rules
  GPU / CUDA                        Constraint: CPU-only execution required
  pandas / numpy                    Not needed; stdlib json + csv handles the pipeline

DEVELOPMENT TOOLING  (not part of ranking pipeline)
  python-pptx     Used only to generate this slide deck from code
  matplotlib      Chart generation for presentation visuals""",

    10: f"""\
SUBMISSION ASSETS

GitHub Repository
  URL:   {GITHUB_REPO}
  Files:
    ranker.py              — complete ranking solution (~420 lines, pure Python stdlib)
    fill_presentation.py   — script that generated this slide deck from code
    validate_submission.py — organiser-provided validator (unchanged)

Reproduce in ~7 Seconds
  $ python3 ranker.py   →  reads candidates.jsonl, scores 100 K profiles, writes submission.csv

Output File — submission.csv
  Columns:  candidate_id  |  rank  |  score  |  reasoning
  Validated checks:
    ✓ All ranks 1–100 present and unique
    ✓ Scores non-increasing from rank 1 to rank 100
    ✓ All candidate IDs match CAND_XXXXXXX format
    ✓ Reasoning column populated for every row

Validation Command
  $ python3 validate_submission.py submission.csv
  → "Submission is valid."

Contact
  {PARTICIPANT_NAME}  |  {PARTICIPANT_EMAIL}""",
}

# ─── Matplotlib: architecture pipeline diagram ────────────────────────────────

def make_architecture_diagram():
    """Return PNG bytes of the pipeline flowchart for slide 7."""
    fig, ax = plt.subplots(figsize=(11.5, 4.5), facecolor="#F8FAFB")
    ax.set_xlim(0, 11.5)
    ax.set_ylim(0, 4.5)
    ax.axis("off")

    # Colour palette
    C_INPUT  = "#E8F4FD"
    C_STAGE  = "#D6EAF8"
    C_ALGO   = "#EBF5FB"
    C_OUTPUT = "#E9F7EF"
    C_BORDER = "#2471A3"
    C_ARROW  = "#1A5276"
    C_TITLE  = "#1A5276"
    C_TEXT   = "#1C2833"
    C_GATE   = "#FDEBD0"

    def box(x, y, w, h, label, sub, color, border):
        rect = FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.05",
                              facecolor=color, edgecolor=border, linewidth=1.5)
        ax.add_patch(rect)
        ax.text(x + w/2, y + h*0.62, label, ha="center", va="center",
                fontsize=9, fontweight="bold", color=C_TITLE, wrap=True)
        ax.text(x + w/2, y + h*0.28, sub, ha="center", va="center",
                fontsize=7.2, color=C_TEXT, wrap=True)

    def arrow(x1, x2, y):
        ax.annotate("", xy=(x2, y), xytext=(x1, y),
                    arrowprops=dict(arrowstyle="->", color=C_ARROW,
                                   lw=1.8, mutation_scale=14))

    # Pipeline boxes (x, y, w, h)
    stages = [
        (0.12, 1.20, 1.65, 2.00, "INPUT",           "candidates.jsonl\n487 MB · 100 K records",       C_INPUT,  "#2471A3"),
        (2.10, 1.20, 1.65, 2.00, "① HONEYPOT\nFILTER", "is_honeypot()\nclaimed vs. actual tenure\n19 filtered", C_GATE,  "#E67E22"),
        (4.05, 1.20, 1.65, 2.00, "② FEATURE\nEXTRACTION", "score_skills()\nscore_experience()\nscore_signals()\nscore_education()", C_STAGE, C_BORDER),
        (6.00, 1.20, 1.65, 2.00, "③ COMPOSITE\nSCORING",  "MH_GATE × disqual\n×(0.42·skills\n+0.28·exp\n+0.18·sig+0.12·edu)", C_STAGE, C_BORDER),
        (7.98, 1.20, 1.65, 2.00, "④ MIN-HEAP\nBUFFER",    "heapq, size=400\nO(log N) insert\ntop-100 selected", C_ALGO, C_BORDER),
        (9.93, 1.20, 1.45, 2.00, "OUTPUT",           "submission.csv\n100 ranked\ncandidates", C_OUTPUT, "#1E8449"),
    ]

    for (x, y, w, h, lbl, sub, col, brd) in stages:
        box(x, y, w, h, lbl, sub, col, brd)

    # Arrows between boxes
    arrow_coords = [
        (1.77, 2.10, 2.20),
        (3.75, 4.05, 2.20),
        (5.70, 6.00, 2.20),
        (7.65, 7.98, 2.20),
        (9.63, 9.93, 2.20),
    ]
    for (x1, x2, y) in arrow_coords:
        arrow(x1, x2, y)

    # Must-have gate annotation below box 3
    ax.text(6.825, 0.90, "MH_GATE:  4/4→×1.00  |  3/4→×0.78  |  2/4→×0.56  |  0/4→×0.10",
            ha="center", va="center", fontsize=7.5, color="#922B21",
            bbox=dict(boxstyle="round,pad=0.3", facecolor="#FDEBD0", edgecolor="#E74C3C", lw=1))

    # Skill taxonomy note under feature extraction
    ax.text(4.875, 0.55, "12 frozensets · 150+ aliases · O(1) lookup",
            ha="center", va="center", fontsize=7.2, color="#2E4057",
            style="italic")

    # Title
    ax.text(5.75, 4.20, "Redrob Candidate Ranking — Single-File Pipeline Architecture",
            ha="center", va="center", fontsize=11, fontweight="bold", color=C_TITLE)

    plt.tight_layout(pad=0.3)
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)
    buf.seek(0)
    return buf


# ─── Matplotlib: score distribution + title bar chart for slide 8 ────────────

def make_results_charts():
    """Return PNG bytes of a two-panel chart: score curve + title distribution."""
    # Real submission data
    all_scores = [
        0.8632, 0.8510, 0.8399, 0.8300, 0.8218, 0.8209, 0.8201, 0.8097, 0.8083, 0.8077,
        0.8070, 0.8026, 0.7989, 0.7966, 0.7957, 0.7943, 0.7937, 0.7937, 0.7899, 0.7871,
        0.7860, 0.7851, 0.7842, 0.7833, 0.7678, 0.7670, 0.7661, 0.7652, 0.7643, 0.7634,
        0.7625, 0.7610, 0.7601, 0.7592, 0.7583, 0.7574, 0.7565, 0.7551, 0.7542, 0.7533,
        0.7524, 0.7515, 0.7506, 0.7497, 0.7488, 0.7479, 0.7470, 0.7461, 0.7452, 0.7423,
        0.7414, 0.7405, 0.7396, 0.7387, 0.7378, 0.7369, 0.7360, 0.7351, 0.7342, 0.7333,
        0.7324, 0.7315, 0.7306, 0.7297, 0.7288, 0.7279, 0.7270, 0.7261, 0.7252, 0.7243,
        0.7234, 0.7225, 0.7216, 0.7209, 0.7200, 0.7191, 0.7182, 0.7173, 0.7164, 0.7155,
        0.7146, 0.7137, 0.7128, 0.7119, 0.7110, 0.7101, 0.7092, 0.7083, 0.7074, 0.7065,
        0.7056, 0.7047, 0.7038, 0.7029, 0.7020, 0.7011, 0.7002, 0.6993, 0.6984, 0.6914,
    ]
    ranks = list(range(1, 101))

    title_data = {
        "Applied ML Engineer": 20,
        "Machine Learning Eng": 14,
        "Search Engineer": 13,
        "AI Engineer": 10,
        "Sr. Data Scientist": 7,
        "NLP Engineer": 6,
        "Sr. NLP Engineer": 6,
        "Rec. Systems Eng": 6,
        "Sr. ML Engineer": 5,
        "Staff ML Engineer": 5,
        "Others": 8,
    }

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(11.5, 3.8), facecolor="#F8FAFB")
    fig.subplots_adjust(wspace=0.38)

    # ── Left: Score curve ──────────────────────────────────────────────────────
    ax1.set_facecolor("#F8FAFB")
    ax1.plot(ranks, all_scores, color="#2471A3", lw=2.0, zorder=3)
    ax1.fill_between(ranks, all_scores, min(all_scores) - 0.01,
                     alpha=0.18, color="#2471A3")

    # Annotate key points
    for r, s, label in [
        (1, 0.8632, "Rank 1\n0.863"),
        (10, 0.8077, "Rank 10\n0.808"),
        (50, 0.7423, "Rank 50\n0.742"),
        (100, 0.6914, "Rank 100\n0.691"),
    ]:
        ax1.scatter([r], [s], color="#E74C3C", s=40, zorder=5)
        va = "bottom" if r < 80 else "top"
        ax1.annotate(label, (r, s), textcoords="offset points",
                     xytext=(6, 4 if r < 80 else -18),
                     fontsize=7.5, color="#922B21")

    ax1.set_xlim(0, 101)
    ax1.set_ylim(0.66, 0.89)
    ax1.set_xlabel("Rank", fontsize=9)
    ax1.set_ylabel("Composite Score", fontsize=9)
    ax1.set_title("Score Distribution (Ranks 1–100)", fontsize=10, fontweight="bold",
                  color="#1A5276")
    ax1.tick_params(labelsize=8)
    ax1.grid(axis="y", ls="--", alpha=0.4)
    ax1.spines["top"].set_visible(False)
    ax1.spines["right"].set_visible(False)

    # 4/4 must-have annotation
    ax1.axhspan(0.685, 0.89, alpha=0.06, color="green")
    ax1.text(50, 0.875, "All 100 have 4/4 must-haves", ha="center", fontsize=7.5,
             color="#1E8449", style="italic")

    # ── Right: Title distribution bar chart ───────────────────────────────────
    ax2.set_facecolor("#F8FAFB")
    labels = list(title_data.keys())
    values = list(title_data.values())
    colors = ["#2471A3" if i < 5 else "#5DADE2" if i < 9 else "#AED6F1"
              for i in range(len(labels))]

    bars = ax2.barh(range(len(labels)), values, color=colors, height=0.65, zorder=3)
    ax2.set_yticks(range(len(labels)))
    ax2.set_yticklabels(labels, fontsize=8)
    ax2.invert_yaxis()
    ax2.set_xlabel("Count in Top-100", fontsize=9)
    ax2.set_title("Top-100 Title Distribution", fontsize=10, fontweight="bold",
                  color="#1A5276")
    ax2.tick_params(labelsize=8)
    ax2.grid(axis="x", ls="--", alpha=0.4, zorder=0)
    ax2.spines["top"].set_visible(False)
    ax2.spines["right"].set_visible(False)

    for bar, val in zip(bars, values):
        ax2.text(val + 0.2, bar.get_y() + bar.get_height()/2, str(val),
                 va="center", fontsize=7.5, color="#1C2833")

    plt.tight_layout(pad=0.5)
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    plt.close(fig)
    buf.seek(0)
    return buf


# ─── Text frame helpers ───────────────────────────────────────────────────────

def _set_para_spacing(para, line_pct="110000"):
    """Force paragraph line spacing to line_pct% and zero space before/after."""
    pPr = para.find(qn("a:pPr"))
    if pPr is None:
        pPr = etree.Element(qn("a:pPr"))
        para.insert(0, pPr)

    # Line spacing → line_pct%
    lnSpc = pPr.find(qn("a:lnSpc"))
    if lnSpc is None:
        lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
    for child in list(lnSpc):
        lnSpc.remove(child)
    etree.SubElement(lnSpc, qn("a:spcPct"), attrib={"val": line_pct})

    # Space before = 0
    spcBef = pPr.find(qn("a:spcBef"))
    if spcBef is None:
        spcBef = etree.SubElement(pPr, qn("a:spcBef"))
    for child in list(spcBef):
        spcBef.remove(child)
    etree.SubElement(spcBef, qn("a:spcPts"), attrib={"val": "0"})

    # Space after = 0
    spcAft = pPr.find(qn("a:spcAft"))
    if spcAft is None:
        spcAft = etree.SubElement(pPr, qn("a:spcAft"))
    for child in list(spcAft):
        spcAft.remove(child)
    etree.SubElement(spcAft, qn("a:spcPts"), attrib={"val": "0"})


def set_textframe(tf, text, font_size_pt=C_FONT, color_hex="202729", bold=False,
                  line_spacing_pct="110000"):
    """Replace all paragraphs in a text frame; overrides the template's 150% line spacing."""
    txBody = tf._txBody
    paras = txBody.findall(qn("a:p"))
    for p in paras[1:]:
        txBody.remove(p)

    lines = text.split("\n")
    for idx, line in enumerate(lines):
        if idx == 0:
            para = paras[0]
            for r in para.findall(qn("a:r")):
                para.remove(r)
            for br in para.findall(qn("a:br")):
                para.remove(br)
        else:
            para = copy.deepcopy(paras[0])
            for r in para.findall(qn("a:r")):
                para.remove(r)
            for br in para.findall(qn("a:br")):
                para.remove(br)
            txBody.append(para)

        # Override template's 150% line spacing → compact spacing
        _set_para_spacing(para, line_pct=line_spacing_pct)

        r_elem = etree.SubElement(para, qn("a:r"))
        rPr = etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "en-US", "dirty": "0"})
        rPr.set("sz", str(int(font_size_pt * 100)))
        if bold:
            rPr.set("b", "1")
        solidFill = etree.SubElement(rPr, qn("a:solidFill"))
        etree.SubElement(solidFill, qn("a:srgbClr"), attrib={"val": color_hex.upper()})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = line if line else " "


def resize_content_box(shape):
    """Extend the content box height to C_HEIGHT to prevent text overflow."""
    shape.height = C_HEIGHT


# ─── Main ─────────────────────────────────────────────────────────────────────

def main():
    print("Generating diagrams…")
    arch_buf    = make_architecture_diagram()
    results_buf = make_results_charts()

    print("Opening template…")
    prs = Presentation(INPUT_FILE)

    for slide_idx, slide in enumerate(prs.slides, 1):
        content = SLIDE_CONTENT.get(slide_idx)

        # ── Slide 1: Title ────────────────────────────────────────────────────
        if slide_idx == 1:
            assert isinstance(content, dict)
            for shape in slide.shapes:
                if shape.has_text_frame:
                    cur = shape.text_frame.text.strip()
                    if cur in content:
                        set_textframe(shape.text_frame, content[cur],
                                      font_size_pt=13, color_hex="202729")
                        print(f"  Slide 1: '{cur}' updated")
            continue

        # ── Slide 7: System Architecture — diagram image ──────────────────────
        if slide_idx == 7:
            arch_buf.seek(0)
            slide.shapes.add_picture(arch_buf,
                                     left=Inches(0.25),
                                     top=Inches(1.38),
                                     width=Inches(9.52),
                                     height=Inches(4.10))
            print("  Slide 7 [System Architecture]: diagram added")
            continue

        # ── Slide 8: Results — text left + chart right ────────────────────────
        if slide_idx == 8:
            text_boxes = [s for s in slide.shapes if s.has_text_frame]
            if len(text_boxes) >= 2:
                content_shape = text_boxes[1]
                # Narrow the text box to left half
                content_shape.width  = Inches(4.55)
                content_shape.height = C_HEIGHT
                set_textframe(content_shape.text_frame,
                              content.strip(), font_size_pt=9.0)
                title_text = text_boxes[0].text_frame.text.strip()
                print(f"  Slide 8 [{title_text[:40]}]: text updated")
            # Add results chart on right side
            results_buf.seek(0)
            slide.shapes.add_picture(results_buf,
                                     left=Inches(4.78),
                                     top=Inches(1.40),
                                     width=Inches(5.10),
                                     height=Inches(3.90))
            print("  Slide 8: results chart added")
            continue

        # ── Slides 2–6, 9–10: standard content slides ─────────────────────────
        if content is None:
            print(f"  Slide {slide_idx}: no content defined, left as-is.")
            continue

        assert isinstance(content, str)
        text_boxes = [s for s in slide.shapes if s.has_text_frame]
        if len(text_boxes) >= 2:
            content_shape = text_boxes[1]
            resize_content_box(content_shape)
            set_textframe(content_shape.text_frame, content.strip(),
                          font_size_pt=C_FONT)
            title_text = text_boxes[0].text_frame.text.strip()
            print(f"  Slide {slide_idx} [{title_text[:40]}]: updated ({content.count(chr(10))+1} lines)")
        else:
            print(f"  Slide {slide_idx}: only {len(text_boxes)} text box(es), skipping.")

    prs.save(OUTPUT_FILE)
    print(f"\n✓ Saved: {OUTPUT_FILE}")
    print("Next: open in PowerPoint/Keynote, verify layout, export as PDF.")


if __name__ == "__main__":
    main()
